from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as ST
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTFILE = "lions-presentation.pptx"

# Warm safari palette
SAND = RGBColor(244, 225, 186)
SAND_LIGHT = RGBColor(252, 244, 224)
CLAY = RGBColor(198, 127, 59)
BROWN = RGBColor(96, 62, 36)
OLIVE = RGBColor(124, 140, 73)
GREEN = RGBColor(99, 120, 62)
WHITE = RGBColor(255, 255, 255)
CHARCOAL = RGBColor(46, 38, 32)
SOFT_RED = RGBColor(181, 83, 63)
SOFT_BLUE = RGBColor(71, 130, 163)


def add_background(slide, top_color=SAND_LIGHT, bottom_color=SAND):
    shapes = slide.shapes
    bg_top = shapes.add_shape(ST.RECTANGLE, 0, 0, Inches(13.333), Inches(3.8))
    bg_top.fill.solid()
    bg_top.fill.fore_color.rgb = top_color
    bg_top.line.fill.background()

    bg_bottom = shapes.add_shape(ST.RECTANGLE, 0, Inches(3.8), Inches(13.333), Inches(3.7))
    bg_bottom.fill.solid()
    bg_bottom.fill.fore_color.rgb = bottom_color
    bg_bottom.line.fill.background()


def set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(8.5), Inches(1.1))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BROWN

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.45), Inches(9), Inches(0.8))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(22)
        sp.font.color.rgb = CHARCOAL


def add_bullets(slide, bullets, left=0.8, top=2.2, width=5.5, height=3.5):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = CHARCOAL
        p.space_after = Pt(10)


def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, SAND_LIGHT, SAND)
    add_title(slide, "Lions", "Kings of the Savannah")

    sun = slide.shapes.add_shape(ST.OVAL, Inches(9.5), Inches(0.9), Inches(1.7), Inches(1.7))
    sun.fill.solid()
    sun.fill.fore_color.rgb = CLAY
    sun.line.fill.background()

    mane = slide.shapes.add_shape(ST.OVAL, Inches(8.4), Inches(2.0), Inches(3.3), Inches(3.3))
    mane.fill.solid()
    mane.fill.fore_color.rgb = CLAY
    mane.line.color.rgb = BROWN
    mane.line.width = Pt(1.5)

    face = slide.shapes.add_shape(ST.OVAL, Inches(9.0), Inches(2.5), Inches(2.1), Inches(2.1))
    face.fill.solid()
    face.fill.fore_color.rgb = SAND
    face.line.color.rgb = BROWN

    for x in (9.5, 10.25):
        eye = slide.shapes.add_shape(ST.OVAL, Inches(x), Inches(3.2), Inches(0.18), Inches(0.18))
        eye.fill.solid()
        eye.fill.fore_color.rgb = BROWN
        eye.line.fill.background()

    nose = slide.shapes.add_shape(ST.ISOSCELES_TRIANGLE, Inches(9.94), Inches(3.52), Inches(0.25), Inches(0.2))
    nose.fill.solid()
    nose.fill.fore_color.rgb = BROWN
    nose.line.fill.background()

    body = slide.shapes.add_shape(ST.ROUNDED_RECTANGLE, Inches(8.9), Inches(4.5), Inches(2.7), Inches(1.2))
    body.fill.solid()
    body.fill.fore_color.rgb = OLIVE
    body.line.fill.background()

    add_bullets(slide, [
        "Scientific name: Panthera leo",
        "Apex predator of African savannah ecosystems",
        "Known for power, teamwork, and iconic roar",
    ], left=0.9, top=2.5, width=6.8, height=3.0)

    set_notes(slide, "Open by framing lions as a symbol of strength and ecological balance. Emphasize that lions are social big cats and top predators.")


def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Habitat & Range")
    add_bullets(slide, [
        "Savannas, grasslands, and open woodlands",
        "Strongholds in eastern and southern Africa",
        "Small remnant population in India (Gir Forest)",
    ])

    map_panel = slide.shapes.add_shape(ST.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.9), Inches(5.8), Inches(4.8))
    map_panel.fill.solid()
    map_panel.fill.fore_color.rgb = WHITE
    map_panel.line.color.rgb = RGBColor(220, 202, 170)

    africa = slide.shapes.add_shape(ST.CLOUD, Inches(7.3), Inches(2.7), Inches(2.9), Inches(2.6))
    africa.fill.solid()
    africa.fill.fore_color.rgb = OLIVE
    africa.line.color.rgb = GREEN

    india = slide.shapes.add_shape(ST.ROUNDED_RECTANGLE, Inches(10.6), Inches(3.3), Inches(1.1), Inches(0.65))
    india.fill.solid()
    india.fill.fore_color.rgb = SOFT_BLUE
    india.line.fill.background()

    for x, y, color in [(8.2, 3.6, SOFT_RED), (8.7, 4.2, SOFT_RED), (10.95, 3.55, CLAY)]:
        m = slide.shapes.add_shape(ST.OVAL, Inches(x), Inches(y), Inches(0.2), Inches(0.2))
        m.fill.solid()
        m.fill.fore_color.rgb = color
        m.line.fill.background()

    label = slide.shapes.add_textbox(Inches(7.1), Inches(5.6), Inches(5.0), Inches(0.8))
    ltf = label.text_frame
    ltf.text = "Map-style view: major range in Africa, small pocket in India"
    ltf.paragraphs[0].font.size = Pt(14)
    ltf.paragraphs[0].font.color.rgb = CHARCOAL
    ltf.paragraphs[0].alignment = PP_ALIGN.CENTER

    set_notes(slide, "Explain how lion range has shrunk over time due to land-use change. Point out Africa as the primary range and India as a small but important population.")


def slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Pride Life")
    add_bullets(slide, [
        "Prides are family groups led by related lionesses",
        "Lionesses coordinate cub care and most hunts",
        "Adult males defend territory and deter rivals",
    ], width=6.0)

    cx, cy = Inches(9.8), Inches(4.1)
    center = slide.shapes.add_shape(ST.OVAL, cx, cy, Inches(1.7), Inches(1.7))
    center.fill.solid()
    center.fill.fore_color.rgb = CLAY
    center.line.color.rgb = BROWN
    ctf = center.text_frame
    ctf.text = "Pride"
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.size = Pt(20)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

    roles = [
        ("Lionesses", Inches(8.0), Inches(2.5), OLIVE),
        ("Cubs", Inches(11.2), Inches(2.5), SOFT_BLUE),
        ("Males", Inches(8.0), Inches(5.5), SOFT_RED),
        ("Territory", Inches(11.2), Inches(5.5), GREEN),
    ]
    for txt, x, y, color in roles:
        box = slide.shapes.add_shape(ST.ROUNDED_RECTANGLE, x, y, Inches(1.8), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        t = box.text_frame
        t.text = txt
        t.paragraphs[0].font.size = Pt(15)
        t.paragraphs[0].font.bold = True
        t.paragraphs[0].font.color.rgb = WHITE
        t.paragraphs[0].alignment = PP_ALIGN.CENTER

        conn = slide.shapes.add_connector(1, x + Inches(0.9), y + Inches(0.45), cx + Inches(0.85), cy + Inches(0.85))
        conn.line.color.rgb = BROWN
        conn.line.width = Pt(1.2)

    set_notes(slide, "Use this infographic to show cooperation inside a pride. Mention that social structure is one reason lions differ from most other big cats.")


def slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Hunting & Diet")
    add_bullets(slide, [
        "Targets include zebra, wildebeest, and antelope",
        "Hunts combine stealth, teamwork, and speed",
        "Lions also scavenge when opportunities appear",
    ], width=5.8)

    steps = [
        ("Stalk", Inches(6.7), OLIVE),
        ("Chase", Inches(8.9), CLAY),
        ("Capture", Inches(11.1), SOFT_RED),
    ]

    for i, (label, x, color) in enumerate(steps):
        card = slide.shapes.add_shape(ST.ROUNDED_RECTANGLE, x, Inches(2.7), Inches(1.9), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        tf = card.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        icon = slide.shapes.add_shape(ST.ISOSCELES_TRIANGLE, x + Inches(0.62), Inches(3.35), Inches(0.65), Inches(0.5))
        icon.fill.solid()
        icon.fill.fore_color.rgb = WHITE
        icon.line.fill.background()

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(ST.RIGHT_ARROW, x + Inches(1.95), Inches(3.45), Inches(0.85), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BROWN
            arrow.line.fill.background()

    diet_band = slide.shapes.add_shape(ST.ROUNDED_RECTANGLE, Inches(6.8), Inches(5.3), Inches(6.0), Inches(1.0))
    diet_band.fill.solid()
    diet_band.fill.fore_color.rgb = WHITE
    diet_band.line.color.rgb = RGBColor(220, 202, 170)
    dt = diet_band.text_frame
    dt.text = "Diet: mostly large herbivores; scavenging helps conserve energy"
    dt.paragraphs[0].font.size = Pt(16)
    dt.paragraphs[0].font.color.rgb = CHARCOAL
    dt.paragraphs[0].alignment = PP_ALIGN.CENTER

    set_notes(slide, "Walk through the sequence from stealth to capture. Clarify that hunts often fail, so energy-efficient behavior like scavenging is important.")


def slide5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Threats & Conservation")

    left = slide.shapes.add_shape(ST.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.9), Inches(4.8))
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(250, 235, 226)
    left.line.color.rgb = SOFT_RED
    lt = left.text_frame
    lt.text = "Threats"
    lt.paragraphs[0].font.bold = True
    lt.paragraphs[0].font.size = Pt(24)
    lt.paragraphs[0].font.color.rgb = SOFT_RED
    for b in ["Habitat fragmentation", "Human-wildlife conflict", "Declining prey base"]:
        p = lt.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(19)
        p.font.color.rgb = CHARCOAL

    right = slide.shapes.add_shape(ST.ROUNDED_RECTANGLE, Inches(6.7), Inches(2.0), Inches(5.9), Inches(4.8))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(231, 243, 224)
    right.line.color.rgb = GREEN
    rt = right.text_frame
    rt.text = "Solutions"
    rt.paragraphs[0].font.bold = True
    rt.paragraphs[0].font.size = Pt(24)
    rt.paragraphs[0].font.color.rgb = GREEN
    for b in ["Protected landscapes", "Community coexistence programs", "Science-led monitoring and policy"]:
        p = rt.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(19)
        p.font.color.rgb = CHARCOAL

    bridge = slide.shapes.add_shape(ST.RIGHT_ARROW, Inches(5.85), Inches(4.05), Inches(1.5), Inches(0.7))
    bridge.fill.solid()
    bridge.fill.fore_color.rgb = CLAY
    bridge.line.fill.background()

    set_notes(slide, "Contrast pressures on lion populations with practical conservation responses. Close by encouraging support for evidence-based conservation programs.")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)

    prs.save(OUTFILE)


if __name__ == "__main__":
    main()
